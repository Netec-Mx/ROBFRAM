"""Reescribe el texto de slides puntuales de un .pptx (por índice de slide y
de shape), preservando la fuente/color/tamaño del primer run de cada párrafo
original — solo cambia las palabras, no el diseño.

Uso (desde código, no CLI): ver Capitulo01/generar.py para un ejemplo.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Pt


def set_text_preservando_estilo(shape, lineas: list[str]) -> None:
    """Reemplaza el texto de un shape línea por línea, reutilizando el
    formato (fuente, tamaño, color, negrita) del primer run del primer
    párrafo original como plantilla para todas las líneas nuevas."""
    tf = shape.text_frame
    primer_parrafo = tf.paragraphs[0]
    if not primer_parrafo.runs:
        # Shape sin runs (placeholder vacío): solo asignar texto plano.
        tf.text = "\n".join(lineas)
        return

    run_plantilla = primer_parrafo.runs[0]
    font = run_plantilla.font
    color_rgb = None
    try:
        color_rgb = font.color.rgb
    except Exception:
        color_rgb = None

    tf.clear()
    for i, linea in enumerate(lineas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = linea
        run.font.size = font.size
        run.font.bold = font.bold
        run.font.name = font.name
        if color_rgb is not None:
            run.font.color.rgb = color_rgb


def aplicar_reemplazos(ruta_pptx: str, reemplazos: dict[int, dict[int, list[str]]]) -> None:
    """reemplazos: {indice_slide: {indice_shape: [lineas_nuevas]}}"""
    prs = Presentation(ruta_pptx)
    for idx_slide, shapes_map in reemplazos.items():
        slide = prs.slides[idx_slide]
        for idx_shape, lineas in shapes_map.items():
            shape = slide.shapes[idx_shape]
            set_text_preservando_estilo(shape, lineas)
    prs.save(ruta_pptx)
    print(f"Texto actualizado en {ruta_pptx}")
